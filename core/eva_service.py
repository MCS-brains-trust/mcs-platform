"""
StatementHub — Eva AI Practice Intelligence Service Layer
==========================================================
Implements the core logic for Eva's three components:
  1. Knowledge Brain — embedding, chunking, semantic search
  2. Chat Interface — context building, RAG retrieval, LLM calls
  3. Finalisation Gate — compliance checks, finding generation

Uses the existing ai_service._call_llm() pattern for LLM calls.
"""
import json
import logging
import math
import os
import uuid
from datetime import timedelta
from decimal import Decimal

from django.conf import settings
from django.utils import timezone

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------
CHUNK_SIZE_TOKENS = 512
CHUNK_OVERLAP_TOKENS = 64
EMBEDDING_MODEL = "text-embedding-3-small"
EMBEDDING_DIMENSIONS = 1536
TOP_K_CHUNKS = 8

# Entity-type to applicable compliance checks mapping
ENTITY_CHECK_MAP = {
    "company": [
        "division_7a", "superannuation", "ato_benchmarks",
        "going_concern", "related_party", "tpar",
    ],
    "trust": [
        "division_7a", "superannuation", "ato_benchmarks",
        "trust_distributions", "going_concern", "related_party", "tpar",
    ],
    "partnership": [
        "superannuation", "ato_benchmarks", "going_concern",
        "related_party", "tpar",
    ],
    "sole_trader": [
        "superannuation", "ato_benchmarks", "going_concern",
    ],
    "smsf": [
        "going_concern", "related_party",
    ],
}


# ---------------------------------------------------------------------------
# Embedding
# ---------------------------------------------------------------------------
def get_embedding(text):
    """
    Generate a vector embedding for the given text using OpenAI's
    text-embedding-3-small model.
    Returns a list of floats (1536 dimensions).
    """
    try:
        from openai import OpenAI
    except ImportError:
        raise ImportError("openai package not installed. Run: pip install openai")

    client = OpenAI()
    response = client.embeddings.create(
        model=EMBEDDING_MODEL,
        input=text,
    )
    return response.data[0].embedding


def cosine_similarity(vec_a, vec_b):
    """Compute cosine similarity between two vectors."""
    dot = sum(a * b for a, b in zip(vec_a, vec_b))
    norm_a = math.sqrt(sum(a * a for a in vec_a))
    norm_b = math.sqrt(sum(b * b for b in vec_b))
    if norm_a == 0 or norm_b == 0:
        return 0.0
    return dot / (norm_a * norm_b)


# ---------------------------------------------------------------------------
# Text Chunking
# ---------------------------------------------------------------------------
def chunk_text(text, chunk_size=CHUNK_SIZE_TOKENS, overlap=CHUNK_OVERLAP_TOKENS):
    """
    Split text into overlapping chunks of approximately chunk_size tokens.
    Uses a simple word-based approximation (1 token ≈ 0.75 words).
    Returns list of dicts: [{text, token_count, chunk_index}]
    """
    words = text.split()
    # Approximate: 1 token ≈ 0.75 words, so chunk_size tokens ≈ chunk_size * 0.75 words
    words_per_chunk = int(chunk_size * 0.75)
    words_overlap = int(overlap * 0.75)

    if not words:
        return []

    chunks = []
    start = 0
    chunk_index = 0

    while start < len(words):
        end = min(start + words_per_chunk, len(words))
        chunk_words = words[start:end]
        chunk_text_str = " ".join(chunk_words)
        # Approximate token count
        approx_tokens = int(len(chunk_words) / 0.75)
        chunks.append({
            "text": chunk_text_str,
            "token_count": approx_tokens,
            "chunk_index": chunk_index,
        })
        chunk_index += 1
        if end >= len(words):
            break
        start = end - words_overlap

    return chunks


# ---------------------------------------------------------------------------
# Document Parsing
# ---------------------------------------------------------------------------
def parse_document(file_path, file_type):
    """
    Extract text from a document file.
    Supports: .docx, .pdf, .txt, .xlsx, .pptx, .msg
    Returns the extracted text as a string.
    """
    file_type = file_type.lower().lstrip(".")

    if file_type == "docx":
        return _parse_docx(file_path)
    elif file_type == "pdf":
        return _parse_pdf(file_path)
    elif file_type == "txt":
        return _parse_txt(file_path)
    elif file_type == "xlsx":
        return _parse_xlsx(file_path)
    elif file_type == "pptx":
        return _parse_pptx(file_path)
    elif file_type == "msg":
        return _parse_msg(file_path)
    else:
        logger.warning(f"Unsupported file type: {file_type}")
        return ""


def _parse_docx(file_path):
    """Extract text from a Word document."""
    try:
        from docx import Document
        doc = Document(file_path)
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    except Exception as e:
        logger.error(f"Failed to parse DOCX {file_path}: {e}")
        return ""


def _parse_pdf(file_path):
    """Extract text from a PDF document."""
    try:
        import pdfplumber
        text_parts = []
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text_parts.append(page_text)
        return "\n".join(text_parts)
    except ImportError:
        # Fallback to PyPDF2
        try:
            from PyPDF2 import PdfReader
            reader = PdfReader(file_path)
            return "\n".join(
                page.extract_text() or "" for page in reader.pages
            )
        except Exception as e:
            logger.error(f"Failed to parse PDF {file_path}: {e}")
            return ""
    except Exception as e:
        logger.error(f"Failed to parse PDF {file_path}: {e}")
        return ""


def _parse_txt(file_path):
    """Read a plain text file."""
    try:
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            return f.read()
    except Exception as e:
        logger.error(f"Failed to parse TXT {file_path}: {e}")
        return ""


def _parse_xlsx(file_path):
    """Extract text from an Excel file."""
    try:
        from openpyxl import load_workbook
        wb = load_workbook(file_path, read_only=True, data_only=True)
        text_parts = []
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                row_text = " | ".join(str(c) for c in row if c is not None)
                if row_text.strip():
                    text_parts.append(row_text)
        return "\n".join(text_parts)
    except Exception as e:
        logger.error(f"Failed to parse XLSX {file_path}: {e}")
        return ""


def _parse_pptx(file_path):
    """Extract text from a PowerPoint file."""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        text_parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            text_parts.append(para.text)
        return "\n".join(text_parts)
    except Exception as e:
        logger.error(f"Failed to parse PPTX {file_path}: {e}")
        return ""


def _parse_msg(file_path):
    """Extract text from an Outlook .msg email file.

    Extracts the subject, sender, date, and body text.
    Also extracts text from any supported attachments
    (PDF, DOCX, TXT, XLSX, PPTX) embedded in the email.
    """
    try:
        import extract_msg
        msg = extract_msg.Message(file_path)
        msg_message = msg.body or ""

        parts = []
        if msg.subject:
            parts.append(f"Subject: {msg.subject}")
        if msg.sender:
            parts.append(f"From: {msg.sender}")
        if msg.date:
            parts.append(f"Date: {msg.date}")
        parts.append("")
        parts.append(msg_message)

        # Extract text from supported attachments
        for attachment in msg.attachments:
            att_name = getattr(attachment, "longFilename", None) or getattr(attachment, "filename", None) or ""
            att_ext = att_name.rsplit(".", 1)[-1].lower() if "." in att_name else ""
            if att_ext in ("pdf", "docx", "txt", "xlsx", "pptx"):
                try:
                    import tempfile as _tempfile
                    with _tempfile.NamedTemporaryFile(suffix=f".{att_ext}", delete=False) as att_tmp:
                        att_tmp.write(attachment.data)
                        att_tmp_path = att_tmp.name
                    att_text = parse_document(att_tmp_path, att_ext)
                    if att_text.strip():
                        parts.append(f"\n--- Attachment: {att_name} ---")
                        parts.append(att_text)
                    import os
                    os.unlink(att_tmp_path)
                except Exception as e:
                    logger.warning(f"Failed to parse .msg attachment {att_name}: {e}")

        msg.close()
        return "\n".join(parts)
    except ImportError:
        logger.error("extract-msg package not installed. Run: pip install extract-msg")
        return ""
    except Exception as e:
        logger.error(f"Failed to parse MSG {file_path}: {e}")
        return ""


# ---------------------------------------------------------------------------
# Knowledge Brain — Semantic Search
# ---------------------------------------------------------------------------
def search_knowledge_brain(query, category_filter=None, top_k=TOP_K_CHUNKS):
    """
    Search the Knowledge Brain for chunks most relevant to the query.
    Uses pgvector cosine distance for efficient similarity search.

    Args:
        query: The search query text
        category_filter: Optional category string to filter documents
        top_k: Number of top results to return

    Returns:
        List of dicts: [{chunk_id, text, document_title, category, similarity}]
    """
    from pgvector.django import CosineDistance

    from core.models import KnowledgeChunk, KnowledgeDocument

    query_embedding = get_embedding(query)

    # Use pgvector native cosine distance search
    chunks_qs = KnowledgeChunk.objects.filter(
        document__is_archived=False,
        document__sync_status=KnowledgeDocument.SyncStatus.SYNCED,
        embedding_vector__isnull=False,
    ).select_related("document").annotate(
        distance=CosineDistance("embedding_vector", query_embedding),
    ).order_by("distance")

    if category_filter:
        chunks_qs = chunks_qs.filter(document__category=category_filter)

    scored = []
    for chunk in chunks_qs[:top_k]:
        scored.append({
            "chunk_id": str(chunk.id),
            "text": chunk.text,
            "document_title": chunk.document.title,
            "document_id": str(chunk.document.id),
            "category": chunk.document.get_category_display(),
            "similarity": 1.0 - (chunk.distance or 0.0),
        })

    return scored


# ---------------------------------------------------------------------------
# Context Building for Eva Chat
# ---------------------------------------------------------------------------
def build_financial_year_context(financial_year):
    """
    Build the structured context payload that accompanies every Eva chat message.
    Includes: entity data, trial balance, journals, officers, Eva findings, amber indicators.
    """
    from core.models import (
        TrialBalanceLine, AdjustingJournal, EntityOfficer,
    )
    from core.models import EvaReview, EvaFinding

    fy = financial_year
    entity = fy.entity

    # Entity info
    entity_context = {
        "entity_name": entity.entity_name,
        "entity_type": entity.get_entity_type_display(),
        "abn": entity.abn or "Not recorded",
        "financial_year": f"{fy.start_date} to {fy.end_date}",
        "year_label": fy.year_label,
        "trading_as": entity.trading_as or "",
        "trustee_name": entity.trustee_name or "",
    }

    # Trial balance (aggregated by account code)
    tb_lines = TrialBalanceLine.objects.filter(
        financial_year=fy
    ).select_related("mapped_line_item").order_by("account_code")

    tb_data = []
    for line in tb_lines:
        net_current = line.debit - line.credit
        net_prior = line.prior_debit - line.prior_credit
        tb_data.append({
            "code": line.account_code,
            "name": line.account_name,
            "current_dr": str(line.debit),
            "current_cr": str(line.credit),
            "prior_dr": str(line.prior_debit),
            "prior_cr": str(line.prior_credit),
            "net_current": str(net_current),
            "net_prior": str(net_prior),
            "mapped_to": line.mapped_line_item.line_item_label if line.mapped_line_item else "Unmapped",
            "is_adjustment": line.is_adjustment,
        })

    # Posted journals
    journals = AdjustingJournal.objects.filter(
        financial_year=fy, status="posted"
    ).order_by("-posted_at")
    journal_data = []
    for j in journals[:20]:  # Last 20 journals
        journal_data.append({
            "reference": j.reference_number or "Draft",
            "description": j.description or "",
            "date": str(j.journal_date) if j.journal_date else "",
            "total_debit": str(j.total_debit),
            "total_credit": str(j.total_credit),
        })

    # Officers (directors, trustees, beneficiaries)
    officers = EntityOfficer.objects.filter(entity=entity)
    officer_data = []
    for o in officers:
        officer_data.append({
            "name": o.full_name,
            "roles": o.roles if hasattr(o, 'roles') and o.roles else [o.get_role_display()] if hasattr(o, 'role') else [],
        })

    # Existing Eva findings (if any review exists)
    eva_findings = []
    latest_review = EvaReview.objects.filter(financial_year=fy).first()
    if latest_review:
        for f in latest_review.findings.filter(domain='financial_statements'):  # Sprint 1b: scope to FS domain
            eva_findings.append({
                "check": f.get_check_name_display(),
                "severity": f.get_severity_display(),
                "title": f.title,
                "status": f.get_status_display(),
            })

    return {
        "entity": entity_context,
        "trial_balance": tb_data,
        "journals": journal_data,
        "officers": officer_data,
        "eva_findings": eva_findings,
    }


# ---------------------------------------------------------------------------
# Eva Chat — Process a message and generate response
# ---------------------------------------------------------------------------
EVA_CHAT_SYSTEM_PROMPT = """You are Eva, the AI Practice Intelligence system for MC & S Pty Ltd, an Australian accounting practice. You have been built from the accumulated expertise of the firm's principals and a comprehensive library of Australian tax and compliance materials.

You are assisting {accountant_name} who is currently working on {entity_name} ({entity_type}) for the financial year ending {year_end_date}.

You have access to: (1) the full trial balance for this financial year with prior year comparatives, (2) all posted journals, (3) entity director/trustee/beneficiary records, (4) retrieved excerpts from the MC & S Knowledge Brain most relevant to the accountant's question.

When answering:
- Ground your response in the retrieved Knowledge Brain documents where applicable. Cite the specific source (e.g., 'Based on TR 2023/1...' or 'Per MC & S Technical Position memo dated March 2024...').
- When the answer relies on data from the current financial year context, state what data was observed (e.g., 'I can see the loan balance is $68,400 as at 30 June 2025...').
- When you do not have sufficient information in context to answer confidently, say so explicitly and suggest where to find the answer.
- You do not make definitive legal or tax advice statements. You provide the relevant framework and recommend the accountant apply their professional judgment.

## Response Format Rules for Chat

You are responding in a chat bubble, not writing a document. Follow these rules for every chat response:

### Layer 1: Direct Answer (Default)
- Lead with the direct answer to the question in the first sentence.
- Maximum 150 words for your initial response.
- Use natural conversational prose. No markdown headers (#, ##). No bullet point lists. No tables.
- Include the key number, key risk, or key action. Cite legislation or account codes inline only where critical to the answer.
- End every response with ONE of:
  (a) An expansion offer: a short question offering more detail on a specific aspect.
  (b) An action offer: an offer to generate a document, run an analysis, draft a workpaper, or perform a platform action.
- Never end a response with a generic summary or restatement. End with a question or an offer.

### Layer 2: Expanded Detail (On Request)
- Only provide expanded detail when the user explicitly asks for it or accepts your expansion offer.
- Maximum 300 words for expanded responses.
- You may use: short inline tables (max 5 rows), bold for key figures, brief numbered steps (max 5).
- Do NOT repeat information from your Layer 1 response. Provide only the new detail requested.
- Continue to end with an action offer if one has not yet been made.

### Layer 3: Action Offers
- Whenever your answer identifies a risk, a compliance issue, a calculation, or a disclosure requirement, offer to perform a concrete platform action.
- Examples of action offers: draft a document, generate a workpaper, calculate a schedule, prepare a disclosure note, create a finding for review.
- Frame action offers as brief questions: 'Want me to draft the loan agreement?' not 'I could potentially prepare a loan agreement document if that would be helpful.'

### Formatting Constraints
- Never use markdown headers (# or ##) in chat responses.
- Never use bullet point lists in your initial (Layer 1) response.
- Replace tables with natural language in Layer 1. Example: instead of a table showing CY vs PY, write 'Revenue dropped from $1.2M to $890K, a 26% decline.'
- If the user asks a yes/no question, answer yes or no in the first word.
- If the user asks for a specific number, state the number in the first sentence.
- Use 'I' naturally. You are Eva, a colleague, not a system.

### What NOT to Change
- Maintain full accuracy. Never sacrifice correctness for brevity.
- Continue citing legislation, account codes, and specific dollar amounts where relevant.
- Continue using Knowledge Brain retrieval for grounded answers.
- Continue cross-referencing conversation history for context.
- If the user's question genuinely requires more than 150 words to answer safely and accurately (e.g., a multi-part compliance question), you may exceed the limit but must still lead with the direct answer and avoid document-style formatting.

### Anti-Patterns — What You Must NOT Do
- Do NOT open with a section header like 'Division 7A Analysis for CST Automation Pty Ltd'. Open with the answer.
- Do NOT include an 'Overview' section before the substance. First sentence = answer.
- Do NOT render evidence or comparisons as multi-column tables. Summarise in prose.
- Do NOT provide a numbered remediation checklist unprompted. State the top 1-2 actions. Offer the full checklist on request.
- Do NOT end with a passive summary like 'In summary, Division 7A compliance requires...'. End with an offer or question.
- Do NOT produce template language blocks (e.g., sample contract clauses) unprompted. Offer to generate the document instead."""

EVA_FINALISATION_SYSTEM_PROMPT = """You are Eva, conducting a formal compliance review of {entity_name} ({entity_type}) for the financial year ending {year_end_date}.

You must check the following applicable compliance areas: {applicable_checks}

For each area, review the financial data provided and determine whether a material compliance issue exists.

Return your findings as a JSON array. Each finding must include:
- check_name: one of {check_names}
- severity: "critical" or "advisory"
- title: brief plain-English title (max 100 chars)
- explanation: plain English, maximum 3 sentences
- recommendation: specific action the accountant should take
- legislation_reference: e.g. 'ITAA 1936 s.109D'
- knowledge_brain_citation: if a firm Knowledge Brain document was cited (otherwise empty string)
- confidence: "high", "medium", or "low"

Rules:
- Only raise findings where the data provides clear evidence of a material issue.
- Set severity to "advisory" for uncertain or ambiguous situations.
- Do not raise findings outside the designated check areas.
- A finding with "low" confidence must be "advisory" regardless of the check type.
- If no material issues are found for any check area, return an empty JSON array: []

Return ONLY the JSON array, no other text."""


def process_eva_chat(financial_year, user, message_text, opus_override=False, interaction_type="general"):
    """
    Process a chat message from the accountant and generate Eva's response.

    Args:
        financial_year: FinancialYear instance
        user: User instance (the accountant)
        message_text: The accountant's question
        opus_override: If True, use Opus model instead of Sonnet

    Returns:
        EvaMessage instance (the assistant response)
    """
    from core.models import EvaConversation, EvaMessage
    from core.models import AuditLog
    from core.ai_service import _call_llm

    # Get or create conversation
    conversation, _ = EvaConversation.objects.get_or_create(
        financial_year=financial_year,
        user=user,
    )

    # Save user message
    user_msg = EvaMessage.objects.create(
        conversation=conversation,
        role=EvaMessage.Role.USER,
        content=message_text,
        interaction_type=interaction_type,
    )
    conversation.message_count = conversation.messages.count()
    conversation.save(update_fields=["message_count", "last_active_at"])

    # Check if this is a trust planning query
    from core.eva_trust_planning import (
        is_trust_planning_query, get_trust_planning_prompt,
        get_or_create_planning_session,
    )
    is_trust_planning = is_trust_planning_query(
        message_text, financial_year.entity.entity_type
    )
    if is_trust_planning:
        interaction_type = "trust_planning"
        user_msg.interaction_type = "trust_planning"
        user_msg.save(update_fields=["interaction_type"])
        # Get or create planning session
        planning_session = get_or_create_planning_session(
            financial_year, conversation, user
        )

    # Build context
    fy_context = build_financial_year_context(financial_year)

    # Search Knowledge Brain for relevant chunks
    retrieved_chunks = []
    try:
        retrieved_chunks = search_knowledge_brain(message_text, top_k=TOP_K_CHUNKS)
    except Exception as e:
        logger.warning(f"Knowledge Brain search failed: {e}")

    # Build the system prompt
    entity = financial_year.entity
    if is_trust_planning:
        from core.eva_trust_planning import TRUST_PLANNING_SYSTEM_PROMPT
        system_prompt = TRUST_PLANNING_SYSTEM_PROMPT.format(
            entity_name=entity.entity_name,
            entity_type=entity.get_entity_type_display(),
            year_end_date=financial_year.end_date.strftime("%d %B %Y"),
        )
    else:
        system_prompt = EVA_CHAT_SYSTEM_PROMPT.format(
            accountant_name=user.get_full_name() or user.email,
            entity_name=entity.entity_name,
            entity_type=entity.get_entity_type_display(),
            year_end_date=financial_year.end_date.strftime("%d %B %Y"),
        )

    # Build the user prompt with context
    kb_context = ""
    if retrieved_chunks:
        kb_context = "\n\n--- KNOWLEDGE BRAIN EXCERPTS ---\n"
        for i, chunk in enumerate(retrieved_chunks, 1):
            kb_context += f"\n[{i}] Source: {chunk['document_title']} ({chunk['category']})\n"
            kb_context += f"{chunk['text']}\n"

    # Get recent conversation history (last 10 exchanges)
    recent_messages = list(
        conversation.messages.order_by("-created_at")[:20]
    )[::-1]  # Reverse to chronological order
    history = ""
    for msg in recent_messages[:-1]:  # Exclude the just-created user message
        role_label = "Accountant" if msg.role == "user" else "Eva"
        history += f"\n{role_label}: {msg.content}\n"

    user_prompt = f"""--- ENTITY CONTEXT ---
{json.dumps(fy_context['entity'], indent=2)}

--- TRIAL BALANCE (Current Year) ---
{json.dumps(fy_context['trial_balance'][:100], indent=2)}
{"[... " + str(len(fy_context['trial_balance']) - 100) + " more accounts ...]" if len(fy_context['trial_balance']) > 100 else ""}

--- POSTED JOURNALS ---
{json.dumps(fy_context['journals'], indent=2)}

--- OFFICERS / DIRECTORS / TRUSTEES ---
{json.dumps(fy_context['officers'], indent=2)}

--- EXISTING EVA FINDINGS ---
{json.dumps(fy_context['eva_findings'], indent=2)}
{kb_context}
{f"--- CONVERSATION HISTORY ---{history}" if history else ""}

--- ACCOUNTANT'S QUESTION ---
{message_text}"""

    # Inject trust planning context if applicable
    if is_trust_planning:
        from core.eva_trust_planning import build_trust_planning_context
        trust_ctx = build_trust_planning_context(financial_year)
        trust_section = f"""\n\n--- TRUST DISTRIBUTION PLANNING CONTEXT ---
=== INCOME STREAMS ===
{json.dumps(trust_ctx['income_summary'], indent=2)}

=== BENEFICIARY PROFILES ===
{json.dumps(trust_ctx['beneficiaries'], indent=2)}

=== COMPLIANCE FLAGS ===
{json.dumps(trust_ctx['compliance_flags'], indent=2)}

=== EXISTING SCENARIOS ===
{json.dumps(trust_ctx.get('existing_scenarios', []), indent=2)}"""
        user_prompt += trust_section

    # Select model tier
    tier = "opus" if opus_override else "sonnet"

    try:
        response_text = _call_llm(
            system_prompt=system_prompt,
            user_prompt=user_prompt,
            tier=tier,
            temperature=0.3,
            max_tokens=4000,
        )
    except Exception as e:
        logger.error(f"Eva chat LLM call failed: {e}")
        response_text = (
            "I apologise — I encountered a technical issue processing your question. "
            "Please try again, or contact the system administrator if this persists."
        )

    # Save assistant message
    chunk_ids = [c["chunk_id"] for c in retrieved_chunks]
    model_label = "opus" if opus_override else "sonnet"
    # Estimate token counts (rough: 1 token ≈ 4 chars)
    prompt_tokens = len(user_prompt) // 4 + len(system_prompt) // 4
    response_tokens = len(response_text) // 4
    assistant_msg = EvaMessage.objects.create(
        conversation=conversation,
        role=EvaMessage.Role.ASSISTANT,
        content=response_text,
        model_used=model_label,
        retrieved_chunk_ids=chunk_ids,
        tokens_used=prompt_tokens + response_tokens,
        token_count_prompt=prompt_tokens,
        token_count_response=response_tokens,
        interaction_type=interaction_type,
    )
    # Link M2M knowledge chunks cited
    if chunk_ids:
        from core.models import KnowledgeChunk
        cited_chunks = KnowledgeChunk.objects.filter(id__in=chunk_ids)
        assistant_msg.knowledge_chunks_cited.set(cited_chunks)
    conversation.message_count = conversation.messages.count()
    conversation.save(update_fields=["message_count", "last_active_at"])

    # Log to activity
    AuditLog.objects.create(
        user=user,
        action=AuditLog.Action.EVA_CHAT,
        description=(
            f"{user.get_full_name() or user.email} asked Eva a question. "
            f"Topic: {message_text[:80]}. Model: {model_label}."
        ),
        affected_object_type="FinancialYear",
        affected_object_id=str(financial_year.pk),
        metadata={
            "message_id": str(user_msg.pk),
            "response_id": str(assistant_msg.pk),
            "model": model_label,
            "kb_docs_cited": [c["document_title"] for c in retrieved_chunks[:3]],
        },
    )

    return assistant_msg


# ---------------------------------------------------------------------------
# Eva Finalisation Gate — Run compliance review
# ---------------------------------------------------------------------------
def run_eva_review(financial_year, user, opus_override=False):
    """
    Run Eva's structured compliance review (Finalisation Gate).

    Args:
        financial_year: FinancialYear instance
        user: User who triggered the review
        opus_override: If True, use Opus model

    Returns:
        EvaReview instance with findings
    """
    from core.models import EvaReview, EvaFinding
    from core.models import AuditLog, FinancialYear
    from core.ai_service import _call_llm

    entity = financial_year.entity
    entity_type = entity.entity_type

    # Determine applicable checks
    applicable_checks = ENTITY_CHECK_MAP.get(entity_type, ["going_concern"])

    # Create the review record
    review = EvaReview.objects.create(
        financial_year=financial_year,
        triggered_by=user,
        model_used="opus" if opus_override else "sonnet",
        opus_override=opus_override,
        status=EvaReview.Status.IN_PROGRESS,
        checks_total=len(applicable_checks),
    )

    # FY stays in_review — no status change needed

    # Log the trigger
    AuditLog.objects.create(
        user=user,
        action=AuditLog.Action.EVA_REVIEW,
        description=(
            f"{user.get_full_name() or user.email} submitted this financial year "
            f"to Eva for compliance review. Model: {'opus' if opus_override else 'sonnet'}."
        ),
        affected_object_type="FinancialYear",
        affected_object_id=str(financial_year.pk),
        metadata={"review_id": str(review.pk), "model": review.model_used},
    )

    # Build context
    fy_context = build_financial_year_context(financial_year)

    # Search Knowledge Brain for compliance-related content
    kb_context = ""
    try:
        for check in applicable_checks:
            chunks = search_knowledge_brain(
                f"Australian {check.replace('_', ' ')} compliance requirements",
                top_k=3,
            )
            if chunks:
                kb_context += f"\n--- Knowledge Brain: {check} ---\n"
                for chunk in chunks:
                    kb_context += f"Source: {chunk['document_title']}\n{chunk['text']}\n\n"
    except Exception as e:
        logger.warning(f"Knowledge Brain search for review failed: {e}")

    # Build check names for the prompt
    check_display = {
        "division_7a": "Division 7A",
        "superannuation": "Superannuation Guarantee",
        "ato_benchmarks": "ATO Industry Benchmarks",
        "trust_distributions": "Trust Distributions",
        "going_concern": "Going Concern",
        "related_party": "Related Party Transactions",
        "tpar": "TPAR Obligations",
    }

    applicable_display = [check_display.get(c, c) for c in applicable_checks]

    system_prompt = EVA_FINALISATION_SYSTEM_PROMPT.format(
        entity_name=entity.entity_name,
        entity_type=entity.get_entity_type_display(),
        year_end_date=financial_year.end_date.strftime("%d %B %Y"),
        applicable_checks=", ".join(applicable_display),
        check_names=json.dumps(applicable_checks),
    )

    user_prompt = f"""--- ENTITY CONTEXT ---
{json.dumps(fy_context['entity'], indent=2)}

--- FULL TRIAL BALANCE ---
{json.dumps(fy_context['trial_balance'], indent=2)}

--- POSTED JOURNALS ---
{json.dumps(fy_context['journals'], indent=2)}

--- OFFICERS / DIRECTORS / TRUSTEES / BENEFICIARIES ---
{json.dumps(fy_context['officers'], indent=2)}
{kb_context}

Please conduct your compliance review now and return the JSON array of findings."""

    tier = "opus" if opus_override else "sonnet"

    try:
        response_text = _call_llm(
            system_prompt=system_prompt,
            user_prompt=user_prompt,
            tier=tier,
            temperature=0.2,
            max_tokens=8000,
        )

        # Parse the JSON response
        # Strip any markdown code fences
        cleaned = response_text.strip()
        if cleaned.startswith("```"):
            cleaned = cleaned.split("\n", 1)[1] if "\n" in cleaned else cleaned[3:]
        if cleaned.endswith("```"):
            cleaned = cleaned[:-3]
        cleaned = cleaned.strip()
        if cleaned.startswith("json"):
            cleaned = cleaned[4:].strip()

        findings_data = json.loads(cleaned)

        if not isinstance(findings_data, list):
            raise ValueError("Expected JSON array of findings")

        # Create finding records
        for fd in findings_data:
            EvaFinding.objects.create(
                eva_review=review,
                check_name=fd.get("check_name", "going_concern"),
                severity=fd.get("severity", "advisory"),
                title=fd.get("title", "Untitled Finding"),
                explanation=fd.get("explanation", ""),
                recommendation=fd.get("recommendation", ""),
                legislation_reference=fd.get("legislation_reference", ""),
                knowledge_brain_citation=fd.get("knowledge_brain_citation", ""),
                confidence=fd.get("confidence", "medium"),
            )

        review.checks_completed = len(applicable_checks)
        review.completed_at = timezone.now()

        fs_findings = review.findings.filter(domain='financial_statements')  # Sprint 1b: scope to FS domain
        if fs_findings.exists():
            review.status = EvaReview.Status.FINDINGS_RAISED
            AuditLog.objects.create(
                user=user,
                action=AuditLog.Action.EVA_REVIEW,
                description=(
                    f"Eva has identified {fs_findings.count()} matter(s) "
                    f"requiring attention before this year can be finalised."
                ),
                affected_object_type="FinancialYear",
                affected_object_id=str(financial_year.pk),
                metadata={
                    "review_id": str(review.pk),
                    "finding_count": fs_findings.count(),
                    "critical_count": fs_findings.filter(severity="critical").count(),
                },
            )
        else:
            review.status = EvaReview.Status.CLEARED
            # FY stays in_review — accountant clicks Finalise manually
            AuditLog.objects.create(
                user=user,
                action=AuditLog.Action.EVA_REVIEW,
                description=(
                    f"Eva completed compliance review. No material issues identified. "
                    f"Model: {review.model_used}. "
                    f"Duration: {(review.completed_at - review.triggered_at).seconds}s."
                ),
                affected_object_type="FinancialYear",
                affected_object_id=str(financial_year.pk),
                metadata={"review_id": str(review.pk)},
            )

        review.save()

    except json.JSONDecodeError as e:
        logger.error(f"Eva review JSON parse error: {e}\nRaw response: {response_text[:500]}")
        review.status = EvaReview.Status.ERROR
        review.error_message = f"Failed to parse Eva's response: {e}"
        review.completed_at = timezone.now()
        review.save()
        # Reset FY status
        financial_year.status = FinancialYear.Status.DRAFT
        financial_year.save(update_fields=["status"])
        AuditLog.objects.create(
            user=user,
            action=AuditLog.Action.EVA_REVIEW,
            description=f"Eva review encountered a parse error. Finalisation blocked. Error: {e}",
            affected_object_type="FinancialYear",
            affected_object_id=str(financial_year.pk),
        )

    except Exception as e:
        logger.error(f"Eva review failed: {e}")
        review.status = EvaReview.Status.ERROR
        review.error_message = str(e)
        review.completed_at = timezone.now()
        review.save()
        # Reset FY status
        financial_year.status = FinancialYear.Status.DRAFT
        financial_year.save(update_fields=["status"])
        AuditLog.objects.create(
            user=user,
            action=AuditLog.Action.EVA_REVIEW,
            description=f"Eva review encountered an API error. Finalisation blocked. Error: {e}",
            affected_object_type="FinancialYear",
            affected_object_id=str(financial_year.pk),
        )

    return review


def resolve_eva_finding(finding, user, resolution_note):
    """
    Mark an Eva finding as addressed with the accountant's resolution note.
    If all findings are now addressed, triggers a re-run.

    Returns:
        (finding, should_rerun: bool)
    """
    from core.models import EvaFinding
    from core.models import AuditLog

    finding.status = EvaFinding.FindingStatus.ADDRESSED
    finding.resolution_note = resolution_note
    finding.resolved_by = user
    finding.resolved_at = timezone.now()
    finding.save()

    # Log
    AuditLog.objects.create(
        user=user,
        action=AuditLog.Action.EVA_FINDING,
        description=(
            f"{user.get_full_name() or user.email} addressed Eva finding: "
            f"{finding.title}. Note: \"{resolution_note[:200]}\""
        ),
        affected_object_type="FinancialYear",
        affected_object_id=str(finding.eva_review.financial_year.pk),
        metadata={
            "finding_id": str(finding.pk),
            "review_id": str(finding.eva_review.pk),
        },
    )

    # Check if all findings are now addressed
    review = finding.eva_review
    open_count = review.findings.filter(domain='financial_statements', status=EvaFinding.FindingStatus.OPEN).count()  # Sprint 1b: scope to FS domain
    should_rerun = (open_count == 0)

    return finding, should_rerun


# ---------------------------------------------------------------------------
# Amber Indicators — Trial Balance variance analysis
# ---------------------------------------------------------------------------
def compute_amber_indicators(financial_year, materiality_pct_revenue=15, materiality_pct_expenses=20):
    """
    Compute amber indicators for all trial balance lines in a financial year.
    Returns a dict keyed by account_code with list of trigger descriptions.

    Trigger conditions:
    1. Significant variance — $ (exceeds section materiality threshold)
    2. Significant variance — % (>15% revenue, >20% expenses)
    3. Account dropped (non-zero PY, zero CY)
    4. Opening balance mismatch
    5. Balance sign change

    Note: "Account added" was removed — a new account appearing is almost always
    intentional and provides no meaningful analytical signal.
    """
    from core.models import TrialBalanceLine
    from collections import defaultdict

    ZERO = Decimal("0")
    indicators = defaultdict(list)

    lines = TrialBalanceLine.objects.filter(
        financial_year=financial_year,
    ).select_related("mapped_line_item")

    # Aggregate lines by account code
    account_data = defaultdict(lambda: {
        "current_dr": ZERO, "current_cr": ZERO,
        "prior_dr": ZERO, "prior_cr": ZERO,
        "opening": ZERO, "prior_closing": ZERO,
        "name": "", "section": "",
    })

    for line in lines:
        code = line.account_code
        d = account_data[code]
        d["current_dr"] += line.debit or ZERO
        d["current_cr"] += line.credit or ZERO
        d["prior_dr"] += line.prior_debit or ZERO
        d["prior_cr"] += line.prior_credit or ZERO
        d["opening"] += line.opening_balance or ZERO
        d["prior_closing"] += line.prior_closing_balance or ZERO
        if not d["name"]:
            d["name"] = line.account_name
        if not d["section"] and line.mapped_line_item:
            d["section"] = (line.mapped_line_item.statement_section or "").lower()

    for code, d in account_data.items():
        net_current = d["current_dr"] - d["current_cr"]
        net_prior = d["prior_dr"] - d["prior_cr"]
        variance = net_current - net_prior
        triggers = []

        # 1 & 2: Significant variance
        if net_prior != ZERO:
            pct = abs(variance / abs(net_prior) * 100)
            section = d["section"]
            threshold = materiality_pct_revenue if ("revenue" in section or "income" in section) else materiality_pct_expenses

            if pct > threshold:
                triggers.append({
                    "type": f"Significant variance ({pct:.1f}%)",
                    "prior": str(net_prior),
                    "current": str(net_current),
                    "movement": str(variance),
                    "pct": f"{pct:.1f}",
                })

        # 3: Account dropped
        if net_prior != ZERO and net_current == ZERO:
            triggers.append({
                "type": "Account dropped",
                "prior": str(net_prior),
                "current": "0.00",
                "movement": str(-net_prior),
                "pct": "-100.0",
            })

        # 4: Opening balance mismatch
        if d["prior_closing"] != ZERO and d["opening"] != d["prior_closing"]:
            triggers.append({
                "type": "Opening balance mismatch",
                "prior": f"PY closing: {d['prior_closing']}",
                "current": f"Opening: {d['opening']}",
                "movement": str(d["opening"] - d["prior_closing"]),
                "pct": "N/A",
            })

        # 5: Balance sign change
        if net_prior != ZERO and net_current != ZERO:
            if (net_prior > ZERO and net_current < ZERO) or (net_prior < ZERO and net_current > ZERO):
                triggers.append({
                    "type": "Balance sign change",
                    "prior": str(net_prior),
                    "current": str(net_current),
                    "movement": str(variance),
                    "pct": "N/A",
                })

        if triggers:
            indicators[code] = triggers

    return dict(indicators)


# ---------------------------------------------------------------------------
# SharePoint Sync — Knowledge Brain document sync
# ---------------------------------------------------------------------------
from core.eva_knowledge import SHAREPOINT_FOLDER_MAP


def sync_knowledge_brain(limit=0):
    """
    Sync the Knowledge Brain from SharePoint via Microsoft Graph API.
    This function is designed to be called as a management command or
    scheduled task (Celery Beat every 2 hours).

    Args:
        limit: Maximum number of documents to process in this run.
               0 = no limit (process all). Use limit=1 with a shell
               loop for memory-constrained servers.

    Requires environment variables:
    - SHAREPOINT_TENANT_ID
    - SHAREPOINT_CLIENT_ID
    - SHAREPOINT_CLIENT_SECRET
    - SHAREPOINT_SITE_URL (e.g. https://tenant.sharepoint.com/sites/SiteName)
    - SHAREPOINT_LIBRARY_NAME (e.g. Eva Knowledge Brain)

    Optionally, you can set SHAREPOINT_SITE_ID and SHAREPOINT_DRIVE_ID directly
    to skip the auto-resolution step.

    Processes documents one at a time with explicit memory cleanup between
    each document to prevent OOM kills on memory-constrained servers.
    """
    import gc
    import time
    import requests as http_requests
    from core.models import KnowledgeDocument, KnowledgeChunk
    from core.models import AuditLog
    import tempfile
    from urllib.parse import urlparse

    tenant_id = os.environ.get("SHAREPOINT_TENANT_ID", "")
    client_id = os.environ.get("SHAREPOINT_CLIENT_ID", "")
    client_secret = os.environ.get("SHAREPOINT_CLIENT_SECRET", "")
    site_url = os.environ.get("SHAREPOINT_SITE_URL", "")
    library_name = os.environ.get("SHAREPOINT_LIBRARY_NAME", "Eva Knowledge Brain")
    site_id = os.environ.get("SHAREPOINT_SITE_ID", "")
    drive_id = os.environ.get("SHAREPOINT_DRIVE_ID", "")

    if not all([tenant_id, client_id, client_secret]):
        logger.error("SharePoint credentials not configured. Set SHAREPOINT_TENANT_ID, SHAREPOINT_CLIENT_ID, SHAREPOINT_CLIENT_SECRET.")
        return {"error": "SharePoint credentials not configured", "added": 0, "updated": 0}

    if not site_id and not site_url:
        logger.error("Set either SHAREPOINT_SITE_ID or SHAREPOINT_SITE_URL.")
        return {"error": "SharePoint site not configured", "added": 0, "updated": 0}

    # Get OAuth2 token
    token_url = f"https://login.microsoftonline.com/{tenant_id}/oauth2/v2.0/token"
    token_data = {
        "grant_type": "client_credentials",
        "client_id": client_id,
        "client_secret": client_secret,
        "scope": "https://graph.microsoft.com/.default",
    }
    token_resp = http_requests.post(token_url, data=token_data)
    token_resp.raise_for_status()
    access_token = token_resp.json()["access_token"]
    headers = {"Authorization": f"Bearer {access_token}"}
    graph_base = "https://graph.microsoft.com/v1.0"

    # Auto-resolve site_id from site_url if not provided
    if not site_id:
        parsed = urlparse(site_url)
        hostname = parsed.hostname  # e.g. mcandscomau.sharepoint.com
        site_path = parsed.path.rstrip("/")  # e.g. /sites/MCS354
        resolve_url = f"{graph_base}/sites/{hostname}:{site_path}"
        logger.info(f"Resolving SharePoint site ID from: {resolve_url}")
        site_resp = http_requests.get(resolve_url, headers=headers)
        site_resp.raise_for_status()
        site_id = site_resp.json()["id"]
        logger.info(f"Resolved site_id: {site_id}")

    # Auto-resolve drive_id from library name if not provided
    if not drive_id:
        drives_url = f"{graph_base}/sites/{site_id}/drives"
        drives_resp = http_requests.get(drives_url, headers=headers)
        drives_resp.raise_for_status()
        for drv in drives_resp.json().get("value", []):
            if drv.get("name") == library_name:
                drive_id = drv["id"]
                break
        if not drive_id:
            logger.error(f"Could not find drive for library '{library_name}'. Available: {[d['name'] for d in drives_resp.json().get('value', [])]}")
            return {"error": f"Drive not found for library '{library_name}'", "added": 0, "updated": 0}
        logger.info(f"Resolved drive_id: {drive_id}")

    stats = {"added": 0, "updated": 0, "skipped": 0, "errors": 0, "total_chunks": 0}
    docs_processed = 0

    # Iterate through each folder in the mapping
    limit_reached = False
    for folder_path, category in SHAREPOINT_FOLDER_MAP.items():
        if limit_reached:
            break
        full_path = folder_path  # Drive root is already the library
        list_url = (
            f"{graph_base}/sites/{site_id}/drives/{drive_id}"
            f"/root:/{full_path}:/children"
        )

        try:
            resp = http_requests.get(list_url, headers=headers)
            if resp.status_code == 404:
                logger.info(f"SharePoint folder not found: {full_path}")
                continue
            resp.raise_for_status()
            items = resp.json().get("value", [])
        except Exception as e:
            logger.error(f"Failed to list SharePoint folder {full_path}: {e}")
            stats["errors"] += 1
            continue

        for item in items:
            if item.get("folder"):
                # Skip Archive subfolders
                if "archive" in item["name"].lower():
                    # Mark any existing docs from this path as archived
                    KnowledgeDocument.objects.filter(
                        sharepoint_path__startswith=f"{full_path}/{item['name']}",
                        is_archived=False,
                    ).update(is_archived=True)
                continue

            item_id = item["id"]
            item_name = item["name"]
            modified_at = item.get("lastModifiedDateTime", "")
            file_ext = item_name.rsplit(".", 1)[-1].lower() if "." in item_name else ""

            if file_ext not in ("docx", "pdf", "txt", "xlsx", "pptx", "msg"):
                continue

            # Check if document already exists and is up to date
            existing = KnowledgeDocument.objects.filter(
                sharepoint_item_id=item_id
            ).first()

            if existing and str(existing.sharepoint_modified_at) == modified_at:
                stats["skipped"] += 1
                continue  # No changes

            # Check if we've hit the document limit for this run
            if limit > 0 and docs_processed >= limit:
                logger.info(f"Reached document limit ({limit}). Stopping.")
                limit_reached = True
                break

            # Download the file
            download_url = (
                f"{graph_base}/sites/{site_id}/drives/{drive_id}"
                f"/items/{item_id}/content"
            )
            try:
                file_resp = http_requests.get(download_url, headers=headers)
                file_resp.raise_for_status()
            except Exception as e:
                logger.error(f"Failed to download {item_name}: {e}")
                stats["errors"] += 1
                continue

            # Save to temp file and parse
            with tempfile.NamedTemporaryFile(suffix=f".{file_ext}", delete=False) as tmp:
                tmp.write(file_resp.content)
                tmp_path = tmp.name

            try:
                text = parse_document(tmp_path, file_ext)
                if not text.strip():
                    logger.warning(f"No text extracted from {item_name}")
                    continue

                # Chunk the text
                chunks = chunk_text(text)

                # Free the raw text and file content immediately
                del text
                file_size = len(file_resp.content)
                del file_resp

                # Create or update the document record
                if existing:
                    doc = existing
                    doc.title = item_name.rsplit(".", 1)[0]
                    doc.sharepoint_modified_at = modified_at
                    doc.file_type = file_ext
                    doc.file_size_bytes = file_size
                    doc.is_archived = False
                    # Delete old chunks
                    doc.chunks.all().delete()
                    stats["updated"] += 1
                else:
                    doc = KnowledgeDocument.objects.create(
                        title=item_name.rsplit(".", 1)[0],
                        category=category,
                        sharepoint_path=f"{full_path}/{item_name}",
                        sharepoint_item_id=item_id,
                        sharepoint_modified_at=modified_at,
                        file_type=file_ext,
                        file_size_bytes=file_size,
                    )
                    stats["added"] += 1

                # Embed and store chunks one at a time with periodic memory cleanup
                num_chunks = len(chunks)
                for i, chunk_data in enumerate(chunks):
                    try:
                        embedding = get_embedding(chunk_data["text"])
                    except Exception as e:
                        logger.error(f"Embedding failed for chunk {chunk_data['chunk_index']} of {item_name}: {e}")
                        embedding = None

                    KnowledgeChunk.objects.create(
                        document=doc,
                        chunk_index=chunk_data["chunk_index"],
                        text=chunk_data["text"],
                        embedding=embedding,
                        token_count=chunk_data["token_count"],
                    )
                    stats["total_chunks"] += 1
                    # Free embedding vector immediately
                    del embedding
                    # Clear the chunk data we just processed
                    chunks[i] = None

                    # Every 5 chunks, sleep briefly and force garbage collection
                    if (i + 1) % 5 == 0:
                        time.sleep(0.5)  # Let the OS reclaim memory
                        gc.collect()
                        from django.db import reset_queries, close_old_connections
                        reset_queries()
                        close_old_connections()

                # Free the chunks list
                del chunks
                gc.collect()

                doc.chunk_count = num_chunks
                doc.sync_status = KnowledgeDocument.SyncStatus.SYNCED
                doc.synced_at = timezone.now()
                doc.save()

                docs_processed += 1

                # Log progress after each document
                logger.info(
                    f"[{docs_processed}] Synced {item_name} "
                    f"({num_chunks} chunks). "
                    f"Total chunks so far: {stats['total_chunks']}"
                )

                # Sleep between documents to let OS stabilise memory
                time.sleep(2)

            except Exception as e:
                logger.error(f"Failed to process {item_name}: {e}")
                if existing:
                    existing.sync_status = KnowledgeDocument.SyncStatus.ERROR
                    existing.sync_error = str(e)
                    existing.save()
                stats["errors"] += 1
            finally:
                try:
                    os.unlink(tmp_path)
                except OSError:
                    pass
                # Force garbage collection after each document
                gc.collect()
                from django.db import reset_queries, close_old_connections
                reset_queries()
                close_old_connections()

    # Log the sync
    try:
        AuditLog.objects.create(
            action=AuditLog.Action.EVA_SYNC,
            description=(
                f"Eva Knowledge Brain synced. Documents added: {stats['added']}. "
                f"Documents updated: {stats['updated']}. "
                f"Total chunks: {stats['total_chunks']}."
            ),
            metadata=stats,
        )
    except Exception:
        pass

    logger.info(f"Knowledge Brain sync complete: {stats}")
    return stats


# ===========================================================================
# Interactive Clarification System
# ===========================================================================

CLARIFICATION_QUESTIONS = {
    "div7a": {
        "question_id": "div7a_borrower_type",
        "question_text": (
            "Eva has flagged a potential Division 7A issue for the loan account "
            "**\"{account_name}\"**. To help Eva assess this accurately, please "
            "confirm: who is the borrower?"
        ),
        "options": [
            {
                "value": "related_company",
                "label": "Related Company (Pty Ltd)",
                "outcome_hint": "dismiss",
                "outcome_message": (
                    "Division 7A does not apply to loans between companies. "
                    "This finding has been dismissed. Eva will remember this "
                    "for future reviews of this entity."
                ),
                "learning_note": (
                    "Borrower confirmed as a related company — Div 7A does not apply. "
                    "Dismiss Div 7A findings for this account in future reviews."
                ),
            },
            {
                "value": "director_shareholder",
                "label": "Director / Shareholder (individual)",
                "outcome_hint": "confirm",
                "outcome_message": (
                    "Loans to directors and shareholders are subject to Division 7A. "
                    "This finding has been confirmed as requiring attention. "
                    "A complying loan agreement or repayment is required before lodgement day."
                ),
                "learning_note": (
                    "Borrower confirmed as a director/shareholder — Div 7A applies. "
                    "Confirm Div 7A findings for this account in future reviews."
                ),
            },
            {
                "value": "associate_spouse",
                "label": "Associate / Spouse / Family Member",
                "outcome_hint": "confirm",
                "outcome_message": (
                    "Loans to associates of shareholders (including spouses and family members) "
                    "are subject to Division 7A under s 109 ITAA 1936. "
                    "This finding has been confirmed as critical."
                ),
                "learning_note": (
                    "Borrower confirmed as an associate/spouse — Div 7A applies. "
                    "Confirm Div 7A findings for this account in future reviews."
                ),
            },
            {
                "value": "trust",
                "label": "Trust (discretionary or unit)",
                "outcome_hint": "needs_review",
                "outcome_message": (
                    "Loans to trusts can attract Division 7A if the trust has a "
                    "shareholder or associate as a beneficiary. This finding has been "
                    "flagged for further review — please check whether a UPE sub-trust "
                    "arrangement or complying loan agreement is in place."
                ),
                "learning_note": (
                    "Borrower confirmed as a trust — Div 7A may apply depending on "
                    "beneficiary composition. Flag for further review in future."
                ),
            },
            {
                "value": "employee_arm_length",
                "label": "Employee (arm's length, not a shareholder/associate)",
                "outcome_hint": "reduce_severity",
                "outcome_message": (
                    "Loans to employees who are not shareholders or associates are "
                    "generally not subject to Division 7A. This finding has been "
                    "reduced to Advisory pending confirmation that the employee has "
                    "no shareholding or associate relationship."
                ),
                "learning_note": (
                    "Borrower confirmed as an arm's length employee — Div 7A unlikely. "
                    "Reduce severity for this account in future reviews."
                ),
            },
        ],
    },
    "related_party": {
        "question_id": "related_party_arm_length",
        "question_text": (
            "Eva has flagged a related party transaction for **\"{account_name}\"**. "
            "Was this transaction conducted at arm's length (i.e., on commercial terms "
            "that would apply between independent parties)?"
        ),
        "options": [
            {
                "value": "arm_length_documented",
                "label": "Yes — arm's length, documented",
                "outcome_hint": "dismiss",
                "outcome_message": (
                    "Arm's length related party transactions are generally acceptable. "
                    "This finding has been dismissed. Ensure the documentation is "
                    "retained on file."
                ),
                "learning_note": "Related party transaction confirmed as arm's length with documentation.",
            },
            {
                "value": "arm_length_undocumented",
                "label": "Yes — arm's length, but not formally documented",
                "outcome_hint": "reduce_severity",
                "outcome_message": (
                    "The transaction appears arm's length but lacks formal documentation. "
                    "This finding has been reduced to Advisory. Consider preparing a "
                    "written agreement or board minute to support the pricing."
                ),
                "learning_note": "Related party transaction arm's length but undocumented — advisory.",
            },
            {
                "value": "not_arm_length",
                "label": "No — not at arm's length",
                "outcome_hint": "confirm",
                "outcome_message": (
                    "Non-arm's length related party transactions may attract ATO scrutiny "
                    "and transfer pricing adjustments. This finding has been confirmed. "
                    "Consider whether the pricing should be adjusted or disclosed."
                ),
                "learning_note": "Related party transaction confirmed as not arm's length.",
            },
        ],
    },
    "super_guarantee": {
        "question_id": "sgc_worker_classification",
        "question_text": (
            "Eva has identified a potential Superannuation Guarantee shortfall for "
            "**\"{account_name}\"**. Is this worker classified as an employee for "
            "superannuation purposes (including contractors paid mainly for labour)?"
        ),
        "options": [
            {
                "value": "employee_sgc_paid",
                "label": "Yes — employee, SGC was paid correctly",
                "outcome_hint": "dismiss",
                "outcome_message": (
                    "SGC has been confirmed as correctly paid for this worker. "
                    "This finding has been dismissed."
                ),
                "learning_note": "Worker confirmed as employee with correct SGC paid.",
            },
            {
                "value": "contractor_no_sgc",
                "label": "Contractor — not subject to SGC",
                "outcome_hint": "reduce_severity",
                "outcome_message": (
                    "If this worker is a genuine contractor (not a deemed employee), "
                    "SGC may not apply. This finding has been reduced to Advisory. "
                    "Ensure the contractor arrangement is documented and does not "
                    "meet the results test criteria under s 12(3) SGAA 1992."
                ),
                "learning_note": "Worker classified as contractor — SGC may not apply.",
            },
            {
                "value": "employee_sgc_underpaid",
                "label": "Yes — employee, SGC may be underpaid",
                "outcome_hint": "confirm",
                "outcome_message": (
                    "SGC underpayment is a critical compliance risk. This finding has "
                    "been confirmed. A Superannuation Guarantee Charge (SGC) statement "
                    "may need to be lodged with the ATO."
                ),
                "learning_note": "Worker confirmed as employee with potential SGC underpayment.",
            },
        ],
    },
    "trust_distribution": {
        "question_id": "trust_resolution_status",
        "question_text": (
            "Eva has flagged a potential issue with trust distribution resolutions "
            "for **\"{account_name}\"**. Has a valid distribution resolution been "
            "passed before 30 June of the relevant income year?"
        ),
        "options": [
            {
                "value": "resolution_in_place",
                "label": "Yes — resolution was passed before 30 June",
                "outcome_hint": "dismiss",
                "outcome_message": (
                    "A valid distribution resolution is in place. This finding has "
                    "been dismissed. Ensure the resolution is retained on file."
                ),
                "learning_note": "Trust distribution resolution confirmed as in place before 30 June.",
            },
            {
                "value": "resolution_late",
                "label": "Resolution passed after 30 June",
                "outcome_hint": "confirm",
                "outcome_message": (
                    "A late distribution resolution may result in the trustee being "
                    "assessed on the entire net income at the top marginal rate. "
                    "This finding has been confirmed as critical."
                ),
                "learning_note": "Trust distribution resolution confirmed as late.",
            },
            {
                "value": "no_resolution",
                "label": "No resolution was passed",
                "outcome_hint": "confirm",
                "outcome_message": (
                    "Without a valid distribution resolution, the trustee may be "
                    "assessed on the entire net income at the top marginal rate under "
                    "s 99A ITAA 1936. This finding has been confirmed as critical."
                ),
                "learning_note": "No trust distribution resolution — critical finding confirmed.",
            },
        ],
    },
}


def get_clarification_question(check_name, account_name=""):
    """
    Return the clarification question definition for a given check_name,
    with the account_name placeholder substituted in the question text.
    Returns None if no clarification question is defined for this check_name.
    """
    defn = CLARIFICATION_QUESTIONS.get(check_name)
    if not defn:
        return None
    question_text = defn["question_text"].replace("{account_name}", account_name or "this account")
    return {
        "question_id": defn["question_id"],
        "question_text": question_text,
        "options": defn["options"],
    }


def _reevaluate_finding(finding, clarification):
    """
    Re-evaluate an EvaFinding based on a submitted clarification answer.
    Updates finding.severity, finding.confidence, and finding.status
    based on the outcome_hint from the selected option.
    Returns a dict with the updated finding state.
    """
    from core.models import EvaFinding

    outcome_hint = clarification.outcome_hint
    outcome_message = clarification.outcome_message

    if outcome_hint == "dismiss":
        finding.status = EvaFinding.FindingStatus.ADDRESSED
        finding.resolution_note = (
            f"[Eva Clarification] {outcome_message}\n"
            f"Answered by: {clarification.answered_by.get_full_name() if clarification.answered_by else 'Unknown'}"
        )
        finding.resolved_by = clarification.answered_by
        finding.resolved_at = timezone.now()
        finding.save(update_fields=["status", "resolution_note", "resolved_by", "resolved_at"])
        clarification.outcome = "dismissed"

        # Remove eva_flags from TB lines now that the finding is addressed
        try:
            from core.eva_engine import untag_tb_lines_for_finding
            fy = finding.eva_review.financial_year
            untag_tb_lines_for_finding(fy, finding.check_name)
        except Exception:
            pass  # non-critical — flag display only

    elif outcome_hint == "confirm":
        if finding.severity != EvaFinding.Severity.CRITICAL:
            finding.severity = EvaFinding.Severity.CRITICAL
        finding.confidence = EvaFinding.Confidence.HIGH
        finding.save(update_fields=["severity", "confidence"])
        clarification.outcome = "confirmed"

    elif outcome_hint == "reduce_severity":
        if finding.severity == EvaFinding.Severity.CRITICAL:
            finding.severity = EvaFinding.Severity.ADVISORY
            finding.save(update_fields=["severity"])
        clarification.outcome = "reduced"

    elif outcome_hint == "needs_review":
        finding.confidence = EvaFinding.Confidence.MEDIUM
        finding.save(update_fields=["confidence"])
        clarification.outcome = "pending"

    clarification.save(update_fields=["outcome"])

    review = finding.eva_review
    open_count = review.findings.filter(domain='financial_statements', status="open").count()  # Sprint 1b: scope to FS domain

    return {
        "finding_id": str(finding.pk),
        "new_status": finding.status,
        "new_severity": finding.severity,
        "new_confidence": finding.confidence,
        "outcome": clarification.outcome,
        "outcome_message": outcome_message,
        "should_clear_review": (open_count == 0),
    }
