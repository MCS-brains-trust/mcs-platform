"""
Eva Knowledge Brain — SharePoint Sync, Document Parsing, Chunking & Embedding

This module handles:
1. SharePoint sync via Microsoft Graph API
2. Document parsing (docx, pdf, txt, xlsx, pptx)
3. Text chunking (~512 tokens with 64-token overlap)
4. Embedding generation via OpenAI text-embedding-3-small (1536 dims)
5. Similarity search for RAG retrieval
"""
import io
import json
import logging
import math
import os
import re
from datetime import datetime

from django.conf import settings
from django.utils import timezone

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Configuration
# ---------------------------------------------------------------------------
CHUNK_SIZE_TOKENS = 512
CHUNK_OVERLAP_TOKENS = 64
EMBEDDING_MODEL = "text-embedding-3-small"
EMBEDDING_DIMENSIONS = 1536
TOP_K_CHUNKS = 8

# Rough token estimation: ~4 chars per token for English text
CHARS_PER_TOKEN = 4


# ---------------------------------------------------------------------------
# SharePoint Folder → Category Map
# ---------------------------------------------------------------------------
# Canonical mapping from SharePoint folder paths to KnowledgeDocument.Category
# values. Each value here MUST exist in core.models.KnowledgeDocument.Category.
SHAREPOINT_FOLDER_MAP = {
    # 01_Accounting_Standards
    "01_Accounting_Standards/AASB_Standards":          "aasb_standards",
    "01_Accounting_Standards/IFRS_References":         "ifrs_references",
    "01_Accounting_Standards/Interpretations":         "interpretations",

    # 02_Legislation (PARENT NAME CHANGED from 02_Tax_Legislation)
    "02_Legislation":                                  "legislation_general",     # NEW (catches parent-level files; currently 0)
    "02_Legislation/ATO_Rulings":                      "ato_rulings",             # was 02_Tax_Legislation/ATO_Rulings
    "02_Legislation/Corporations":                     "corporations_act",        # NEW (7 files, latest 2026-03-02)
    "02_Legislation/FBT":                              "fbt",                     # was 02_Tax_Legislation/FBT
    "02_Legislation/GST_BAS":                          "gst_bas",                 # was 02_Tax_Legislation/GST_BAS
    "02_Legislation/Income_Tax":                       "income_tax",              # was 02_Tax_Legislation/Income_Tax
    "02_Legislation/Partnerships":                     "legislation_partnerships",# NEW (3 files, latest 2026-03-02)
    "02_Legislation/Superannuation":                   "superannuation",          # was 02_Tax_Legislation/Superannuation
    "02_Legislation/Tax_Administration":               "tax_administration",      # NEW (4 files, latest 2026-02-27)
    "02_Legislation/Trusts":                           "legislation_trusts",      # NEW (3 files, latest 2026-03-02)

    # 03_Firm_Policies
    "03_Firm_Policies":                                "firm_policies_general",   # NEW (catches 9 parent-level files)
    "03_Firm_Policies/Engagement_Letters":             "firm_engagement",
    "03_Firm_Policies/Quality_Control":                "firm_quality",
    "03_Firm_Policies/Review_Checklists":              "firm_checklists",
    "03_Firm_Policies/Style_Guides":                   "firm_style",

    # 04_Disclosure_Templates
    "04_Disclosure_Templates":                         "disclosure_general",      # NEW (catches 5 parent-level files)
    "04_Disclosure_Templates/Companies":               "disclosure_companies",
    "04_Disclosure_Templates/Partnerships":            "disclosure_partnerships",
    "04_Disclosure_Templates/SMSF":                    "disclosure_smsf",
    "04_Disclosure_Templates/Sole_Traders":            "disclosure_sole_traders",
    "04_Disclosure_Templates/Trusts":                  "disclosure_trusts",

    # 05_Industry_Guides (unchanged)
    "05_Industry_Guides/Construction":                 "industry_construction",
    "05_Industry_Guides/Hospitality":                  "industry_hospitality",
    "05_Industry_Guides/Medical":                      "industry_medical",
    "05_Industry_Guides/Not_For_Profit":               "industry_nfp",
    "05_Industry_Guides/Professional_Services":        "industry_professional",
    "05_Industry_Guides/Retail":                       "industry_retail",

    # 06_Training_Materials
    "06_Training_Materials":                           "training_general",        # NEW (catches 2 parent-level files)
    "06_Training_Materials/Eva_User_Guide":            "training_eva",
    "06_Training_Materials/Onboarding":                "training_onboarding",
    "06_Training_Materials/StatementHub_Procedures":   "training_statementhub",

    # 07_Benchmarks
    "07_Benchmarks/ATO_Benchmarks":                    "ato_benchmarks",
    "07_Benchmarks/Financial_Ratios":                  "financial_ratios",
    "07_Benchmarks/Industry_Benchmarks":               "industry_benchmarks",

    # 08_ATO_Rulings_and_Alerts (NEW top-level)
    "08_ATO_Rulings_and_Alerts":                       "ato_rulings_and_alerts",  # NEW (currently 0 files)

    # 09_Precedents_and_Technical_Positions (NEW top-level — 28 Knowledge Shop files)
    "09_Precedents_and_Technical_Positions":           "firm_precedents",         # NEW

    # 10_Legal_Templates (NEW top-level — 46 files)
    "10_Legal_Templates":                              "legal_general",           # NEW (currently 0 parent-level)
    "10_Legal_Templates/Approved_Templates":           "legal_approved",          # NEW (6 files)
    "10_Legal_Templates/Converted_Templates":          "legal_converted",         # NEW (8 files)
    "10_Legal_Templates/InFinity_Source":              "legal_infinity",          # NEW (32 files)

    # 11_CoWorker_Tax_Return (NEW top-level)
    "11_CoWorker_Tax_Return":                          "tax_return_workpapers",   # NEW (currently 0 files)
}


# ---------------------------------------------------------------------------
# SharePoint Sync via Microsoft Graph API
# ---------------------------------------------------------------------------
def _get_graph_token():
    """Obtain an access token for Microsoft Graph API using client credentials."""
    import requests

    tenant_id = os.environ.get("SHAREPOINT_TENANT_ID", "")
    client_id = os.environ.get("SHAREPOINT_CLIENT_ID", "")
    client_secret = os.environ.get("SHAREPOINT_CLIENT_SECRET", "")

    if not all([tenant_id, client_id, client_secret]):
        raise ValueError(
            "SharePoint credentials not configured. Set SHAREPOINT_TENANT_ID, "
            "SHAREPOINT_CLIENT_ID, SHAREPOINT_CLIENT_SECRET environment variables."
        )

    url = f"https://login.microsoftonline.com/{tenant_id}/oauth2/v2.0/token"
    data = {
        "grant_type": "client_credentials",
        "client_id": client_id,
        "client_secret": client_secret,
        "scope": "https://graph.microsoft.com/.default",
    }
    resp = requests.post(url, data=data, timeout=30)
    resp.raise_for_status()
    return resp.json()["access_token"]


def _get_graph_headers():
    """Get authorization headers for Graph API calls."""
    token = _get_graph_token()
    return {"Authorization": f"Bearer {token}", "Accept": "application/json"}


def sync_sharepoint_library(site_id=None, drive_id=None, folder_path=""):
    """
    Sync documents from a SharePoint document library.

    Args:
        site_id: SharePoint site ID (from env if not provided)
        drive_id: Drive ID (from env if not provided)
        folder_path: Subfolder path within the drive (optional)

    Returns:
        dict with counts: {"synced": N, "skipped": N, "errors": N}
    """
    import requests
    from core.models import KnowledgeDocument

    site_id = site_id or os.environ.get("SHAREPOINT_SITE_ID", "")
    drive_id = drive_id or os.environ.get("SHAREPOINT_DRIVE_ID", "")

    if not site_id or not drive_id:
        raise ValueError("SHAREPOINT_SITE_ID and SHAREPOINT_DRIVE_ID must be set.")

    headers = _get_graph_headers()

    # List items in the drive/folder
    if folder_path:
        url = f"https://graph.microsoft.com/v1.0/drives/{drive_id}/root:/{folder_path}:/children"
    else:
        url = f"https://graph.microsoft.com/v1.0/drives/{drive_id}/root/children"

    counts = {"synced": 0, "skipped": 0, "errors": 0}
    supported_extensions = {".docx", ".pdf", ".txt", ".xlsx", ".pptx"}
    max_pages = 200  # Safety limit to prevent infinite pagination loops
    seen_urls = set()

    while url:
        resp = requests.get(url, headers=headers, timeout=30)
        resp.raise_for_status()
        data = resp.json()

        for item in data.get("value", []):
            # Skip folders — recurse into them
            if "folder" in item:
                sub_path = f"{folder_path}/{item['name']}" if folder_path else item["name"]
                sub_counts = sync_sharepoint_library(site_id, drive_id, sub_path)
                for k in counts:
                    counts[k] += sub_counts[k]
                continue

            # Skip unsupported file types
            name = item.get("name", "")
            ext = os.path.splitext(name)[1].lower()
            if ext not in supported_extensions:
                counts["skipped"] += 1
                continue

            try:
                sp_item_id = item["id"]
                sp_modified = item.get("lastModifiedDateTime", "")
                sp_path = f"{folder_path}/{name}" if folder_path else name
                file_size = item.get("size", 0)

                # Resolve folder path → category via SHAREPOINT_FOLDER_MAP
                # (case-insensitive prefix match). Falls back to firm_procedures
                # with a warning if no folder in the map matches.
                resolved_category = None
                folder_lower = (folder_path or "").lower()
                for mapped_folder, mapped_category in SHAREPOINT_FOLDER_MAP.items():
                    if folder_lower.startswith(mapped_folder.lower()):
                        resolved_category = mapped_category
                        break
                if resolved_category is None:
                    logger.warning(
                        "SharePoint folder '%s' did not match any entry in "
                        "SHAREPOINT_FOLDER_MAP — defaulting to firm_procedures",
                        folder_path,
                    )
                    resolved_category = "firm_procedures"

                # Check if document already exists and is up to date
                existing = KnowledgeDocument.objects.filter(
                    sharepoint_item_id=sp_item_id
                ).first()

                if existing and existing.sync_status == "synced":
                    if sp_modified and existing.sharepoint_modified_at:
                        remote_dt = datetime.fromisoformat(
                            sp_modified.replace("Z", "+00:00")
                        )
                        if remote_dt <= existing.sharepoint_modified_at:
                            counts["skipped"] += 1
                            continue

                # Download the file content
                download_url = item.get("@microsoft.graph.downloadUrl", "")
                if not download_url:
                    download_url = (
                        f"https://graph.microsoft.com/v1.0/drives/{drive_id}"
                        f"/items/{sp_item_id}/content"
                    )
                file_resp = requests.get(
                    download_url, headers=headers, timeout=120
                )
                file_resp.raise_for_status()
                file_bytes = file_resp.content

                # Create or update the document record
                doc, created = KnowledgeDocument.objects.update_or_create(
                    sharepoint_item_id=sp_item_id,
                    defaults={
                        "title": os.path.splitext(name)[0],
                        "category": resolved_category,
                        "sharepoint_path": sp_path,
                        "sharepoint_modified_at": (
                            datetime.fromisoformat(sp_modified.replace("Z", "+00:00"))
                            if sp_modified else None
                        ),
                        "file_type": ext.lstrip("."),
                        "file_size_bytes": file_size,
                        "sync_status": "pending",
                    },
                )

                # Parse, chunk, and embed
                text = parse_document(file_bytes, ext)
                chunks = chunk_text(text)
                embed_and_store_chunks(doc, chunks)

                doc.sync_status = "synced"
                doc.synced_at = timezone.now()
                doc.chunk_count = len(chunks)
                doc.save(update_fields=[
                    "sync_status", "synced_at", "chunk_count"
                ])

                counts["synced"] += 1
                logger.info(f"Synced: {sp_path} ({len(chunks)} chunks)")

            except Exception as e:
                counts["errors"] += 1
                logger.error(f"Error syncing {name}: {e}")
                if existing:
                    existing.sync_status = "error"
                    existing.save(update_fields=["sync_status"])

        # Pagination — with infinite loop protection
        next_url = data.get("@odata.nextLink")
        if next_url and next_url in seen_urls:
            logger.warning("SharePoint pagination loop detected — same nextLink URL returned twice. Stopping.")
            break
        if next_url:
            seen_urls.add(next_url)
        max_pages -= 1
        if max_pages <= 0:
            logger.warning("SharePoint pagination safety limit reached (200 pages). Stopping.")
            break
        url = next_url

    return counts


# ---------------------------------------------------------------------------
# Document Parsing
# ---------------------------------------------------------------------------
def parse_document(file_bytes, extension):
    """
    Parse a document's bytes into plain text.

    Args:
        file_bytes: Raw bytes of the file
        extension: File extension (e.g. '.docx', '.pdf')

    Returns:
        Extracted text as a string
    """
    ext = extension.lower().lstrip(".")

    if ext == "txt":
        return file_bytes.decode("utf-8", errors="replace")

    elif ext == "docx":
        return _parse_docx(file_bytes)

    elif ext == "pdf":
        return _parse_pdf(file_bytes)

    elif ext == "xlsx":
        return _parse_xlsx(file_bytes)

    elif ext == "pptx":
        return _parse_pptx(file_bytes)

    else:
        logger.warning(f"Unsupported file type: {ext}")
        return ""


def _parse_docx(file_bytes):
    """Extract text from a DOCX file."""
    try:
        from docx import Document

        doc = Document(io.BytesIO(file_bytes))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        # Also extract text from tables
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    paragraphs.append(" | ".join(cells))
        return "\n\n".join(paragraphs)
    except Exception as e:
        logger.error(f"DOCX parse error: {e}")
        return ""


def _parse_pdf(file_bytes):
    """Extract text from a PDF file."""
    try:
        import subprocess
        import tempfile

        with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name

        result = subprocess.run(
            ["pdftotext", "-layout", tmp_path, "-"],
            capture_output=True, text=True, timeout=60,
        )
        os.unlink(tmp_path)
        return result.stdout.strip()
    except Exception as e:
        logger.error(f"PDF parse error: {e}")
        return ""


def _parse_xlsx(file_bytes):
    """Extract text from an XLSX file."""
    try:
        from openpyxl import load_workbook

        wb = load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
        lines = []
        for ws in wb.worksheets:
            lines.append(f"=== Sheet: {ws.title} ===")
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                if any(c.strip() for c in cells):
                    lines.append(" | ".join(cells))
        return "\n".join(lines)
    except Exception as e:
        logger.error(f"XLSX parse error: {e}")
        return ""


def _parse_pptx(file_bytes):
    """Extract text from a PPTX file."""
    try:
        from pptx import Presentation

        prs = Presentation(io.BytesIO(file_bytes))
        lines = []
        for i, slide in enumerate(prs.slides, 1):
            lines.append(f"=== Slide {i} ===")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            lines.append(text)
        return "\n\n".join(lines)
    except ImportError:
        logger.warning("python-pptx not installed. PPTX parsing unavailable.")
        return ""
    except Exception as e:
        logger.error(f"PPTX parse error: {e}")
        return ""


# ---------------------------------------------------------------------------
# Text Chunking
# ---------------------------------------------------------------------------
def chunk_text(text, chunk_size=CHUNK_SIZE_TOKENS, overlap=CHUNK_OVERLAP_TOKENS):
    """
    Split text into overlapping chunks of approximately `chunk_size` tokens.

    Uses a simple character-based approximation (4 chars ≈ 1 token) and
    tries to break on paragraph/sentence boundaries.

    Returns:
        List of (chunk_index, chunk_text) tuples
    """
    if not text or not text.strip():
        return []

    # Split into paragraphs first
    paragraphs = re.split(r"\n\s*\n", text)
    paragraphs = [p.strip() for p in paragraphs if p.strip()]

    chunk_char_size = chunk_size * CHARS_PER_TOKEN
    overlap_chars = overlap * CHARS_PER_TOKEN

    chunks = []
    current_chunk = ""
    chunk_idx = 0

    for para in paragraphs:
        # If adding this paragraph would exceed the chunk size
        if len(current_chunk) + len(para) + 2 > chunk_char_size and current_chunk:
            chunks.append((chunk_idx, current_chunk.strip()))
            chunk_idx += 1
            # Keep overlap from end of current chunk
            if overlap_chars > 0 and len(current_chunk) > overlap_chars:
                current_chunk = current_chunk[-overlap_chars:]
            else:
                current_chunk = ""

        if current_chunk:
            current_chunk += "\n\n" + para
        else:
            current_chunk = para

        # If a single paragraph is very long, split it further
        while len(current_chunk) > chunk_char_size:
            # Try to split on sentence boundary
            split_point = chunk_char_size
            for sep in [". ", ".\n", "! ", "? ", "; ", ", "]:
                idx = current_chunk.rfind(sep, 0, chunk_char_size)
                if idx > chunk_char_size // 2:
                    split_point = idx + len(sep)
                    break

            chunks.append((chunk_idx, current_chunk[:split_point].strip()))
            chunk_idx += 1
            current_chunk = current_chunk[split_point - overlap_chars:].strip()

    # Don't forget the last chunk
    if current_chunk.strip():
        chunks.append((chunk_idx, current_chunk.strip()))

    return chunks


# ---------------------------------------------------------------------------
# Embedding Generation & Storage
# ---------------------------------------------------------------------------
def _get_embeddings(texts):
    """
    Generate embeddings for a list of texts using OpenAI API.

    Returns:
        List of embedding vectors (each a list of floats)
    """
    try:
        from openai import OpenAI

        client = OpenAI()
        response = client.embeddings.create(
            model=EMBEDDING_MODEL,
            input=texts,
            dimensions=EMBEDDING_DIMENSIONS,
        )
        return [item.embedding for item in response.data]
    except Exception as e:
        logger.error(f"Embedding generation error: {e}")
        return [[] for _ in texts]


def embed_and_store_chunks(document, chunks):
    """
    Generate embeddings for chunks and store them in the database.

    Args:
        document: KnowledgeDocument instance
        chunks: List of (chunk_index, chunk_text) tuples
    """
    from core.models import KnowledgeChunk

    # Delete existing chunks for this document (re-sync)
    document.chunks.all().delete()

    if not chunks:
        return

    # Batch embed (max 2048 per API call)
    batch_size = 100
    all_embeddings = []

    for i in range(0, len(chunks), batch_size):
        batch_texts = [text for _, text in chunks[i:i + batch_size]]
        embeddings = _get_embeddings(batch_texts)
        all_embeddings.extend(embeddings)

    # Create chunk records — write to both JSON and pgvector fields
    chunk_objects = []
    for (idx, text), embedding in zip(chunks, all_embeddings):
        token_count = max(1, len(text) // CHARS_PER_TOKEN)
        chunk_objects.append(
            KnowledgeChunk(
                document=document,
                chunk_index=idx,
                text=text,
                embedding=embedding,
                embedding_vector=embedding if embedding else None,
                token_count=token_count,
            )
        )

    KnowledgeChunk.objects.bulk_create(chunk_objects)
    logger.info(f"Stored {len(chunk_objects)} chunks for '{document.title}'")


# ---------------------------------------------------------------------------
# RAG Retrieval — Cosine Similarity Search
# ---------------------------------------------------------------------------
def _cosine_similarity(vec_a, vec_b):
    """Compute cosine similarity between two vectors."""
    if not vec_a or not vec_b or len(vec_a) != len(vec_b):
        return 0.0

    dot_product = sum(a * b for a, b in zip(vec_a, vec_b))
    norm_a = math.sqrt(sum(a * a for a in vec_a))
    norm_b = math.sqrt(sum(b * b for b in vec_b))

    if norm_a == 0 or norm_b == 0:
        return 0.0

    return dot_product / (norm_a * norm_b)


def retrieve_relevant_chunks(query, top_k=TOP_K_CHUNKS, category_filter=None):
    """
    Retrieve the top-k most relevant knowledge chunks for a query
    using pgvector cosine distance for efficient similarity search.

    Args:
        query: The search query text
        top_k: Number of chunks to return (default 8)
        category_filter: Optional list of KnowledgeDocument.Category values

    Returns:
        List of dicts: [{"chunk": KnowledgeChunk, "score": float, "document_title": str}]
    """
    from pgvector.django import CosineDistance

    from core.models import KnowledgeChunk

    # Generate query embedding
    query_embeddings = _get_embeddings([query])
    if not query_embeddings or not query_embeddings[0]:
        logger.warning("Failed to generate query embedding")
        return []

    query_vec = query_embeddings[0]

    # Build queryset with pgvector cosine distance
    qs = KnowledgeChunk.objects.select_related("document").filter(
        document__is_archived=False,
        document__sync_status="synced",
        embedding_vector__isnull=False,
    ).annotate(
        distance=CosineDistance("embedding_vector", query_vec),
    ).order_by("distance")

    if category_filter:
        qs = qs.filter(document__category__in=category_filter)

    results = []
    for chunk in qs[:top_k]:
        # CosineDistance returns distance (0 = identical); convert to similarity
        score = 1.0 - (chunk.distance or 0.0)
        results.append({
            "chunk": chunk,
            "score": score,
            "document_title": chunk.document.title,
            "category": chunk.document.get_category_display(),
            "chunk_id": str(chunk.pk),
        })

    return results


def format_rag_context(chunks_with_scores):
    """
    Format retrieved chunks into a context string for the LLM prompt.

    Args:
        chunks_with_scores: Output from retrieve_relevant_chunks()

    Returns:
        Formatted context string
    """
    if not chunks_with_scores:
        return ""

    lines = ["=== KNOWLEDGE BRAIN CONTEXT ===\n"]
    for i, item in enumerate(chunks_with_scores, 1):
        chunk = item["chunk"]
        lines.append(
            f"[Source {i}: {item['document_title']} "
            f"({item['category']}) — Relevance: {item['score']:.2f}]\n"
            f"{chunk.text}\n"
        )

    return "\n".join(lines)


# ---------------------------------------------------------------------------
# Manual Document Upload (non-SharePoint)
# ---------------------------------------------------------------------------
def upload_knowledge_document(file_bytes, filename, category="firm_procedures"):
    """
    Upload a document directly to the Knowledge Brain (without SharePoint).

    Args:
        file_bytes: Raw file bytes
        filename: Original filename
        category: KnowledgeDocument.Category value

    Returns:
        KnowledgeDocument instance
    """
    from core.models import KnowledgeDocument

    ext = os.path.splitext(filename)[1].lower()
    title = os.path.splitext(filename)[0]

    doc = KnowledgeDocument.objects.create(
        title=title,
        category=category,
        file_type=ext.lstrip("."),
        file_size_bytes=len(file_bytes),
        sync_status="pending",
    )

    try:
        text = parse_document(file_bytes, ext)
        chunks = chunk_text(text)
        embed_and_store_chunks(doc, chunks)

        doc.sync_status = "synced"
        doc.synced_at = timezone.now()
        doc.chunk_count = len(chunks)
        doc.save(update_fields=["sync_status", "synced_at", "chunk_count"])

        logger.info(f"Uploaded: {filename} ({len(chunks)} chunks)")
    except Exception as e:
        doc.sync_status = "error"
        doc.save(update_fields=["sync_status"])
        logger.error(f"Upload error for {filename}: {e}")
        raise

    return doc
